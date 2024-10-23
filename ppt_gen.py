from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
import io
import re


def sanitize_string(input_str):
    sanitized = re.sub(r"[^A-Za-z0-9_.-]", "", input_str)
    sanitized = re.sub(r"\.{2,}", ".", sanitized)
    sanitized = re.sub(r"^[^A-Za-z0-9]+", "", sanitized)
    sanitized = re.sub(r"[^A-Za-z0-9]+$", "", sanitized)
    sanitized = sanitized[:63] if len(sanitized) > 63 else sanitized.ljust(3, "_")
    return sanitized

def replace_text(prs, search_text, replace_text):

    # Create a case-insensitive pattern for search_text
    pattern = re.compile(re.escape(search_text), re.IGNORECASE)

    # Loop through all slides in the presentation
    for slide in prs.slides:
        # Loop through all shapes on the slide
        for shape in slide.shapes:
            # Check if the shape has a text frame (contains text)
            if not shape.has_text_frame:
                continue
            
            # Loop through paragraphs and runs within the text frame
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Replace text using the case-insensitive pattern
                    print(run.text)
                    if re.search(pattern, run.text):
                        run.text = re.sub(pattern, replace_text, run.text)
    
    return prs

def ppt_gen(slide_data):
    ppt = Presentation("template_2.pptx")

    ppt = replace_text(ppt, "[title]", slide_data[0][0])
    ppt = replace_text(ppt, "[subtitle]", slide_data[0][1])

    i = 1
    for content in slide_data[1]:
        ppt = replace_text(ppt, f"[slide_title_{i}]", content)
        i += 1

    i = 1
    # Content Slides
    for curr_slide_data in slide_data[2:]:
        ppt = replace_text(ppt, f"[slide_top_title_{i}]", curr_slide_data[0])

        j = 1
        for content in curr_slide_data[1:]:
            ppt = replace_text(ppt, f"[slide_{i}_point_{j}]", content)
            j += 1
        
        i += 1


    # Save the presentation
    ppt_stream = io.BytesIO()
    ppt.save(ppt_stream)
    ppt_stream.seek(0)

    return ppt_stream
